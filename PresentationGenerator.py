from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Cm
from pptx.util import Pt
from PIL import Image
import os

fig_dir = r"C:\Users\mgorbun\Documents\Presentations\fig"
out_file = r"C:\Users\mgorbun\Documents\Presentations\auto.pptx"

prs = Presentation()
blank = prs.slide_layouts[6]

for fname in sorted(os.listdir(fig_dir)):
    if fname.lower().endswith((".png",".jpg",".jpeg",".tif",".tiff",".bmp")):

        slide = prs.slides.add_slide(blank)
        path = os.path.join(fig_dir, fname)

        # желаемая ширина
        width = Cm(23)

        # вычисляем пропорциональную высоту
        img = Image.open(path)
        w_px, h_px = img.size
        height = width * h_px / w_px

        # центрирование
        left = (prs.slide_width - width) / 1.1
        top  = (prs.slide_height - height) / 1.1

        slide.shapes.add_picture(path, left, top, width=width)
        
        txBox = slide.shapes.add_textbox(Cm(0), Cm(0), Cm(1), Cm(1))  # размер пока не важен

        tf = txBox.text_frame
        tf.text = fname[:3]
        
        p = tf.paragraphs[0]
        p.font.size = Pt(40)
        
        tf.word_wrap = False
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

prs.save(out_file)

print('That\'s it')